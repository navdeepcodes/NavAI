"""Document reader — extracts text from PDF, DOCX, PPTX, and plain text files.

All processing is local. No cloud services.
"""
from __future__ import annotations

import csv
import io
import json
from pathlib import Path

from tools.filesystem.path_utils import resolve_path

MAX_TEXT_CHARS = 12_000
MAX_FILE_BYTES = 50_000_000


class DocumentUnreadable(Exception):
    """Raised when a file exists but no usable text could be taken from it.

    Distinct from FileNotFoundError because the remedy differs: a missing file
    means the path is wrong, an unreadable one means the format or the file
    itself is the problem, and the caller should be told which.
    """


# Proportion of non-printable characters above which a "text" extraction is
# really binary. Generous, so that documents with unusual glyphs still read.
_BINARY_RATIO = 0.15


def _looks_binary(text: str) -> bool:
    sample = text[:4000]
    if not sample:
        return False
    unprintable = sum(
        1 for ch in sample
        if ch not in "\n\r\t" and (ord(ch) < 32 or ord(ch) == 127)
    )
    return unprintable / len(sample) > _BINARY_RATIO


def read_document(path: str) -> str:
    file = resolve_path(path)

    if not file.exists():
        raise FileNotFoundError(f"File not found: {file}")

    size = file.stat().st_size
    if size > MAX_FILE_BYTES:
        return f"This file is too large ({size / 1_000_000:.1f} MB). Maximum supported size is {MAX_FILE_BYTES / 1_000_000:.0f} MB."

    suffix = file.suffix.lower()

    if suffix == ".pdf":
        text = _read_pdf(file)
    elif suffix == ".docx":
        text = _read_docx(file)
    elif suffix == ".pptx":
        text = _read_pptx(file)
    elif suffix == ".csv":
        text = _read_csv(file)
    elif suffix == ".json":
        text = _read_json(file)
    elif suffix in (".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
                     ".html", ".css", ".yaml", ".yml", ".toml", ".ini",
                     ".cfg", ".sh", ".bash", ".zsh", ".rb", ".go", ".rs",
                     ".java", ".kt", ".swift", ".c", ".cpp", ".h", ".hpp",
                     ".sql", ".xml", ".env", ".gitignore", ".dockerfile",
                     ".makefile", ".r", ".m", ".lua", ".pl", ".php",
                     ".scala", ".ex", ".exs", ".hs", ".clj", ".dart",
                     ".vue", ".svelte", ""):
        text = _read_text(file)
    else:
        text = _try_text(file)

    # Nothing extracted is a failure, not a document whose contents happen to
    # be an apology. These used to be returned as ordinary text with a success
    # status, so a model received "Could not extract text from this .csv file"
    # in the slot where the file's contents belong -- and had no way to tell
    # that apart from a document that really says that.
    if not text or not text.strip():
        if suffix == ".pdf":
            raise DocumentUnreadable(
                f"{file.name} appears to be image-based or has no extractable "
                "text layer. Scanned pages need OCR, which Mike cannot do yet. "
                "If you can see the text on screen, see_screen may work."
            )
        raise DocumentUnreadable(
            f"No text could be extracted from {file.name} ({suffix or 'no extension'})."
        )

    # Binary content dressed up as text is worse than no content: it gives the
    # model something to pattern-match against. A PNG was being returned as a
    # successful read whose body began with the literal bytes of the header.
    if _looks_binary(text):
        raise DocumentUnreadable(
            f"{file.name} is not a text document — it looks like binary data "
            f"({suffix or 'no extension'}). Mike can read PDF, DOCX, PPTX, CSV, "
            "JSON and plain-text formats. For an image, use see_screen instead."
        )

    if len(text) > MAX_TEXT_CHARS:
        truncated = text[:MAX_TEXT_CHARS]
        return (
            f"{truncated}\n\n"
            f"--- Truncated: showing first {MAX_TEXT_CHARS:,} of {len(text):,} characters. "
            f"The full document is {len(text):,} characters long. ---"
        )

    return text


def _read_pdf(file: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(file))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text.strip()}")

    return "\n\n".join(pages)


def _read_docx(file: Path) -> str:
    from docx import Document

    doc = Document(str(file))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _read_pptx(file: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(file))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))

    return "\n\n".join(slides)


def _read_csv(file: Path) -> str:
    text = file.read_text(encoding="utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for i, row in enumerate(reader):
        if i >= 100:
            rows.append(f"... ({i} rows shown of more)")
            break
        rows.append(" | ".join(row))
    return "\n".join(rows)


def _read_json(file: Path) -> str:
    text = file.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(text)
        formatted = json.dumps(data, indent=2, ensure_ascii=False)
        return formatted
    except json.JSONDecodeError:
        return text


def _read_text(file: Path) -> str:
    return file.read_text(encoding="utf-8", errors="replace")


def _try_text(file: Path) -> str:
    try:
        return file.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return ""
